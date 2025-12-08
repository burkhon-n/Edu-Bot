"""
Text extraction utilities from PDF, DOCX, and PPTX files.
"""

from pathlib import Path
from typing import Optional
import PyPDF2
from docx import Document
from pptx import Presentation


MAX_TEXT_LENGTH = 50000  # Maximum characters to extract


def extract_text_from_pdf(filepath: Path) -> str:
    """Extract text from PDF file."""
    try:
        text = []
        with open(filepath, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            
            # Limit pages to prevent excessive processing
            max_pages = min(num_pages, 50)
            
            for page_num in range(max_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
                
                # Stop if we have enough text
                current_length = len(''.join(text))
                if current_length >= MAX_TEXT_LENGTH:
                    break
        
        full_text = '\n'.join(text)
        return full_text[:MAX_TEXT_LENGTH]
    
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""


def extract_text_from_docx(filepath: Path) -> str:
    """Extract text from DOCX file."""
    try:
        doc = Document(filepath)
        text = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
            
            # Stop if we have enough text
            current_length = len('\n'.join(text))
            if current_length >= MAX_TEXT_LENGTH:
                break
        
        full_text = '\n'.join(text)
        return full_text[:MAX_TEXT_LENGTH]
    
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""


def extract_text_from_pptx(filepath: Path) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(filepath)
        text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
            
            # Stop if we have enough text
            current_length = len('\n'.join(text))
            if current_length >= MAX_TEXT_LENGTH:
                break
        
        full_text = '\n'.join(text)
        return full_text[:MAX_TEXT_LENGTH]
    
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""


def extract_text_from_file(filepath: Path) -> Optional[str]:
    """
    Extract text from supported file formats.
    
    Args:
        filepath: Path to file
        
    Returns:
        Extracted text or None if format not supported
    """
    suffix = filepath.suffix.lower()
    
    if suffix == '.pdf':
        return extract_text_from_pdf(filepath)
    elif suffix in ['.docx', '.doc']:
        return extract_text_from_docx(filepath)
    elif suffix in ['.pptx', '.ppt']:
        return extract_text_from_pptx(filepath)
    else:
        print(f"Unsupported file format: {suffix}")
        return None


def truncate_text_smart(text: str, max_length: int = 10000) -> str:
    """
    Truncate text intelligently, trying to break at sentence boundaries.
    
    Args:
        text: Text to truncate
        max_length: Maximum length
        
    Returns:
        Truncated text
    """
    if len(text) <= max_length:
        return text
    
    # Try to find last sentence ending before max_length
    truncated = text[:max_length]
    
    # Look for sentence endings
    for separator in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
        last_idx = truncated.rfind(separator)
        if last_idx > max_length * 0.8:  # At least 80% of desired length
            return text[:last_idx + 1]
    
    # Fallback: just truncate at max_length
    return truncated
